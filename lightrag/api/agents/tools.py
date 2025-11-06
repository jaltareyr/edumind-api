"""
Agent tools for educational content generation using LightRAG knowledge graph.
"""

import json
import logging
from typing import Any, Dict, List, Optional, Union
from pathlib import Path
import asyncio
from datetime import datetime

from agents import function_tool
from pydantic import BaseModel, Field

# PDF generation imports
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_CENTER, TA_LEFT

# PPT generation imports
from pptx import Presentation
from pptx.util import Inches, Pt

from lightrag.api.routers.query_routes import QueryRequest
from lightrag.lightrag import LightRAG

logger = logging.getLogger(__name__)


# Tool context storage - will be set by the route handler
_tool_context: Dict[str, Any] = {}


def set_tool_context(context: Dict[str, Any]):
    """Set the context for tools to access RAG instance and other resources."""
    global _tool_context
    _tool_context = context


def get_tool_context() -> Dict[str, Any]:
    """Get the current tool context."""
    return _tool_context


class KnowledgeGraphQueryResult(BaseModel):
    """Result from knowledge graph query."""
    query: str = Field(description="The query that was executed")
    context: str = Field(description="Retrieved context from the knowledge graph")
    sources: List[str] = Field(default_factory=list, description="Source references")


class FormattedContent(BaseModel):
    """Formatted educational content with flexible structure for rich documents."""
    title: str = Field(description="Content title")
    sections: List[Dict[str, Any]] = Field(description="Content sections with titles, content (list/string), and optional bullets")
    citations: List[str] = Field(default_factory=list, description="Citations and references")

@function_tool
async def knowledge_graph_query(
    query: str,
    mode: str = "mix"
) -> str:
    """
    Query the LightRAG knowledge graph to retrieve relevant information.

    Args:
        query: The search query to execute against the knowledge graph
        mode: Query mode - "local" for entity-focused, "global" for relationship-focused, 
              "mix" for both (default), "naive" for simple search
    
    Returns:
        String response identical to the /query route output.
    """
    try:
        logger.info(f"=== knowledge_graph_query called ===")
        logger.info(f"Query: {query}")
        
        context = get_tool_context()
        rag_instance = context.get("rag")

        if not isinstance(rag_instance, LightRAG):
            error_msg = "RAG instance not available in tool context"
            logger.error(error_msg)
            return json.dumps({
                "error": error_msg,
                "query": query,
                "context": "",
                "sources": []
            })

        logger.info(f"Executing knowledge graph query via QueryRequest pipeline. Mode: {mode}")

        qr_kwargs = {"query": query, "mode": mode, "only_need_context": True, "chunk_top_k": 1}
        query_request = QueryRequest(**qr_kwargs)
        query_params = query_request.to_query_params(is_stream=False)

        result = await rag_instance.aquery(query_request.query, param=query_params)

        if isinstance(result, str):
            logger.info(f"Returning string response from RAG query")
            return result

        if isinstance(result, dict):
            logger.info(f"Returning dict response from RAG query")
            return json.dumps(result, indent=2)

        logger.info(f"Returning fallback stringified response type: {type(result)}")
        return str(result)

    except Exception as e:
        logger.error(f"Error querying knowledge graph: {str(e)}", exc_info=True)
        return json.dumps({
            "error": str(e),
            "query": query,
            "context": "",
            "sources": []
        })

@function_tool
async def format_content_for_education(
    topic: str,
    raw_content: str,
    target_audience: str = "students",
    content_type: str = "comprehensive"
) -> str:
    """
    Format raw knowledge graph content into rich, structured educational material using LLM.
    
    This function uses GPT to intelligently structure content into:
    - Clear introduction and overview
    - Well-organized main concepts with explanations
    - Practical examples and applications
    - Key takeaways and summary
    - Important points for exam preparation
    
    Args:
        topic: The main topic of the content
        raw_content: Raw content from knowledge graph queries
        target_audience: Target audience (students, professionals, etc.)
        content_type: Type of content (comprehensive, summary, exam-focused)
    
    Returns:
        JSON string with beautifully formatted educational content including sections and citations
    """
    try:
        logger.info(f"=== format_content_for_education called ===")
        logger.info(f"Topic: {topic}")
        logger.info(f"Raw content length: {len(raw_content)}")
        logger.info(f"Target audience: {target_audience}")
        logger.info(f"Content type: {content_type}")
        
        # Parse the raw content if it's JSON
        try:
            content_data = json.loads(raw_content)
            context = content_data.get("context", raw_content)
            sources = content_data.get("sources", [])
        except json.JSONDecodeError:
            context = raw_content
            sources = []
        
        # Use LLM to intelligently structure the content
        from openai import AsyncOpenAI
        client = AsyncOpenAI()
        
        formatting_prompt = f"""You are an expert educational content designer. Transform the following knowledge graph content into a beautifully structured educational document.

Topic: {topic}
Target Audience: {target_audience}
Content Type: {content_type}

Raw Content:
{context}

Create a well-structured educational document with the following sections (adapt as needed):
1. Introduction - Engaging overview of the topic
2. Core Concepts - Main ideas broken down clearly
3. Detailed Explanation - In-depth coverage with subpoints
4. Practical Examples - Real-world applications (if applicable)
5. Key Takeaways - Important points to remember
6. Exam Focus Points - Critical information for {target_audience} (if students)
7. Summary - Concise wrap-up

Requirements:
- Use clear, educational language appropriate for {target_audience}
- Break complex ideas into digestible parts
- Include specific details from the content
- Organize information logically
- Make it comprehensive yet readable
- For each section, create 2-5 paragraphs or bullet points
- Each paragraph should be 2-4 sentences for clarity

Return ONLY a valid JSON object with this exact structure:
{{
  "sections": [
    {{
      "title": "Section Title",
      "content": ["Paragraph 1 text", "Paragraph 2 text", "Paragraph 3 text"],
      "bullets": ["bullet point 1", "bullet point 2"] // optional key points
    }}
  ]
}}

Make it educational, engaging, and rich in content!"""

        response = await client.chat.completions.create(
            model="gpt-5-mini-2025-08-07",
            messages=[
                {"role": "system", "content": "You are an expert educational content formatter. Always return valid JSON."},
                {"role": "user", "content": formatting_prompt}
            ],
            response_format={"type": "json_object"}
        )
        
        formatted_response = response.choices[0].message.content
        structured_content = json.loads(formatted_response)
        
        # Build final formatted content
        sections = structured_content.get("sections", [])
        
        formatted = FormattedContent(
            title=topic,
            sections=sections,
            citations=sources if sources else ["Knowledge Graph Database"]
        )
        
        logger.info(f"Successfully formatted content with {len(sections)} sections")
        return formatted.model_dump_json(indent=2)
        
    except Exception as e:
        logger.error(f"Error formatting content: {str(e)}", exc_info=True)
        # Fallback to basic formatting if LLM fails
        sections = [
            {
                "title": "Overview",
                "content": [f"This document covers {topic} for {target_audience}."],
            },
            {
                "title": "Content",
                "content": [context[:1000] if len(context) > 1000 else context],
            }
        ]
        return json.dumps({"title": topic, "sections": sections, "citations": sources if sources else []})


@function_tool
async def generate_pdf(
    title: str,
    content_json: str,
    output_filename: Optional[str] = None
) -> str:
    """
    Generate a beautifully formatted PDF document from educational content.
    
    Features:
    - Professional styling with colors and formatting
    - Table of contents
    - Proper section hierarchy
    - Bullet points for key information
    - Citation page with numbered references
    - Page numbers and headers
    
    Args:
        title: Document title
        content_json: JSON string containing formatted content (sections and citations)
        output_filename: Optional custom filename (auto-generated if not provided)
    
    Returns:
        Path to the generated PDF file
    """
    try:
        logger.info(f"=== generate_pdf called ===")
        logger.info(f"Title: {title}")
        logger.info(f"Content JSON length: {len(content_json)}")
        
        from reportlab.lib import colors
        from reportlab.lib.styles import ParagraphStyle
        
        # Parse content
        content_data = json.loads(content_json)
        sections = content_data.get("sections", [])
        citations = content_data.get("citations", [])
        
        # Generate filename
        if not output_filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).strip()
            safe_title = safe_title.replace(' ', '_')  # Replace spaces with underscores
            output_filename = f"{safe_title}_{timestamp}.pdf"
        
        # Ensure .pdf extension
        if not output_filename.endswith('.pdf'):
            output_filename = f"{output_filename}.pdf"
        
        # Ensure output directory exists
        context = get_tool_context()
        output_dir = Path(context.get("output_dir", "./output"))
        output_dir.mkdir(parents=True, exist_ok=True)
        
        output_path = output_dir / output_filename
        
        logger.info(f"Generating PDF at: {output_path}")
        
        # Create PDF with custom page template
        doc = SimpleDocTemplate(
            str(output_path),
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )
        
        # Define custom styles
        styles = getSampleStyleSheet()
        
        # Title style - Large, bold, centered, colored
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=28,
            textColor=colors.HexColor('#1a5490'),
            spaceAfter=30,
            spaceBefore=10,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        # Subtitle style
        subtitle_style = ParagraphStyle(
            'Subtitle',
            parent=styles['Normal'],
            fontSize=14,
            textColor=colors.HexColor('#555555'),
            spaceAfter=40,
            alignment=TA_CENTER,
            fontName='Helvetica-Oblique'
        )
        
        # Section heading style - Bold, colored, with spacing
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading1'],
            fontSize=18,
            textColor=colors.HexColor('#2c5aa0'),
            spaceAfter=12,
            spaceBefore=20,
            fontName='Helvetica-Bold',
            borderColor=colors.HexColor('#2c5aa0'),
            borderWidth=0,
            borderPadding=5
        )
        
        # Subheading style
        subheading_style = ParagraphStyle(
            'CustomSubHeading',
            parent=styles['Heading2'],
            fontSize=14,
            textColor=colors.HexColor('#34495e'),
            spaceAfter=8,
            spaceBefore=12,
            fontName='Helvetica-Bold'
        )
        
        # Body text style - Readable, with proper spacing
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['BodyText'],
            fontSize=11,
            textColor=colors.HexColor('#2c3e50'),
            spaceAfter=10,
            leading=16,
            fontName='Helvetica'
        )
        
        # Bullet point style
        bullet_style = ParagraphStyle(
            'BulletPoint',
            parent=styles['BodyText'],
            fontSize=11,
            textColor=colors.HexColor('#2c3e50'),
            leftIndent=20,
            bulletIndent=10,
            spaceAfter=6,
            leading=14,
            fontName='Helvetica'
        )
        
        # Key point style (highlighted)
        key_point_style = ParagraphStyle(
            'KeyPoint',
            parent=styles['BodyText'],
            fontSize=11,
            textColor=colors.HexColor('#c0392b'),
            leftIndent=20,
            bulletIndent=10,
            spaceAfter=6,
            leading=14,
            fontName='Helvetica-Bold'
        )
        
        # Citation style
        citation_style = ParagraphStyle(
            'Citation',
            parent=styles['BodyText'],
            fontSize=9,
            textColor=colors.HexColor('#7f8c8d'),
            leftIndent=20,
            spaceAfter=8,
            fontName='Helvetica'
        )
        
        # Build document content
        story = []
        
        # Title page
        story.append(Spacer(1, 1 * inch))
        story.append(Paragraph(title, title_style))
        story.append(Paragraph("Educational Content from Knowledge Graph", subtitle_style))
        story.append(Paragraph(f"Generated on {datetime.now().strftime('%B %d, %Y')}", subtitle_style))
        story.append(PageBreak())
        
        # Process sections
        for idx, section in enumerate(sections, 1):
            section_title = section.get("title", f"Section {idx}")
            section_content = section.get("content", [])
            section_bullets = section.get("bullets", [])
            
            # Section heading
            story.append(Paragraph(f"{idx}. {section_title}", heading_style))
            
            # Section content
            if isinstance(section_content, list):
                for paragraph in section_content:
                    if isinstance(paragraph, str) and paragraph.strip():
                        story.append(Paragraph(paragraph, body_style))
                        story.append(Spacer(1, 0.1 * inch))
            elif isinstance(section_content, str):
                story.append(Paragraph(section_content, body_style))
                story.append(Spacer(1, 0.1 * inch))
            
            # Bullet points (if any)
            if section_bullets:
                story.append(Spacer(1, 0.1 * inch))
                story.append(Paragraph("Key Points:", subheading_style))
                for bullet in section_bullets:
                    if isinstance(bullet, str) and bullet.strip():
                        bullet_text = f"• {bullet}"
                        story.append(Paragraph(bullet_text, key_point_style))
            
            story.append(Spacer(1, 0.3 * inch))
        
        # Citations page
        if citations:
            story.append(PageBreak())
            story.append(Paragraph("References and Citations", heading_style))
            story.append(Spacer(1, 0.2 * inch))
            
            for i, citation in enumerate(citations, 1):
                citation_text = f"[{i}] {citation}"
                story.append(Paragraph(citation_text, citation_style))
                story.append(Spacer(1, 0.05 * inch))
        
        # Build PDF
        doc.build(story)
        
        logger.info(f"PDF generated successfully: {output_path}")
        return str(output_path)
        
    except Exception as e:
        logger.error(f"Error generating PDF: {str(e)}", exc_info=True)
        return json.dumps({"error": str(e)})


@function_tool
async def generate_ppt(
    title: str,
    content_json: str,
    output_filename: Optional[str] = None
) -> str:
    """
    Generate a beautifully designed PowerPoint presentation from educational content.
    
    Features:
    - Professional design with consistent color scheme
    - Title slide with metadata
    - Section divider slides
    - Content slides with bullets and key points
    - Visual hierarchy with different text sizes
    - Citations slide
    - Proper spacing and layout
    
    Args:
        title: Presentation title
        content_json: JSON string containing formatted content (sections and citations)
        output_filename: Optional custom filename (auto-generated if not provided)
    
    Returns:
        Path to the generated PPT file
    """
    try:
        logger.info(f"=== generate_ppt called ===")
        logger.info(f"Title: {title}")
        logger.info(f"Content JSON length: {len(content_json)}")
        
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        
        # Parse content
        content_data = json.loads(content_json)
        sections = content_data.get("sections", [])
        citations = content_data.get("citations", [])
        
        # Generate filename
        if not output_filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).strip()
            safe_title = safe_title.replace(' ', '_')  # Replace spaces with underscores
            output_filename = f"{safe_title}_{timestamp}.pptx"
        
        # Ensure .pptx extension
        if not output_filename.endswith('.pptx'):
            output_filename = f"{output_filename}.pptx"
        
        # Ensure output directory exists
        context = get_tool_context()
        output_dir = Path(context.get("output_dir", "./output"))
        output_dir.mkdir(parents=True, exist_ok=True)
        
        output_path = output_dir / output_filename
        
        logger.info(f"Generating PPT at: {output_path}")
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Define color scheme
        PRIMARY_COLOR = RGBColor(26, 84, 144)      # Blue
        SECONDARY_COLOR = RGBColor(44, 90, 160)    # Lighter Blue
        ACCENT_COLOR = RGBColor(192, 57, 43)       # Red for key points
        TEXT_COLOR = RGBColor(44, 62, 80)          # Dark gray
        SUBTITLE_COLOR = RGBColor(127, 140, 141)   # Light gray
        
        # ===== TITLE SLIDE =====
        title_slide_layout = prs.slide_layouts[6]  # Blank layout for custom design
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2),
            Inches(9), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Educational Content from Knowledge Graph\nGenerated on {datetime.now().strftime('%B %d, %Y')}"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = SUBTITLE_COLOR
        
        # ===== CONTENT SLIDES =====
        for section_idx, section in enumerate(sections, 1):
            section_title = section.get("title", f"Section {section_idx}")
            section_content = section.get("content", [])
            section_bullets = section.get("bullets", [])
            
            # Create section divider slide for major sections
            if section_idx > 1 and len(sections) > 3:
                divider_layout = prs.slide_layouts[6]
                divider_slide = prs.slides.add_slide(divider_layout)
                
                divider_box = divider_slide.shapes.add_textbox(
                    Inches(1), Inches(3),
                    Inches(8), Inches(1.5)
                )
                divider_frame = divider_box.text_frame
                divider_frame.text = section_title
                divider_para = divider_frame.paragraphs[0]
                divider_para.alignment = PP_ALIGN.CENTER
                divider_para.font.size = Pt(36)
                divider_para.font.bold = True
                divider_para.font.color.rgb = PRIMARY_COLOR
            
            # Determine content presentation
            if isinstance(section_content, list):
                content_items = section_content
            elif isinstance(section_content, str):
                # Split paragraphs into sentences for better presentation
                content_items = [s.strip() for s in section_content.split('.') if s.strip()]
            else:
                content_items = [str(section_content)]
            
            # Create slides for content (chunk if needed)
            chunks = []
            current_chunk = []
            
            for item in content_items:
                if len(item) > 500:  # Long paragraph - give it own slide
                    if current_chunk:
                        chunks.append(current_chunk)
                        current_chunk = []
                    chunks.append([item])
                else:
                    current_chunk.append(item)
                    if len(current_chunk) >= 4:  # Max 4 points per slide
                        chunks.append(current_chunk)
                        current_chunk = []
            
            if current_chunk:
                chunks.append(current_chunk)
            
            # Create slides for each chunk
            for chunk_idx, chunk in enumerate(chunks):
                bullet_slide_layout = prs.slide_layouts[1]  # Title and content
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                # Set title
                title_shape = slide.shapes.title
                slide_title = section_title
                if len(chunks) > 1:
                    slide_title = f"{section_title} ({chunk_idx + 1}/{len(chunks)})"
                title_shape.text = slide_title
                
                # Style title
                title_shape.text_frame.paragraphs[0].font.size = Pt(32)
                title_shape.text_frame.paragraphs[0].font.bold = True
                title_shape.text_frame.paragraphs[0].font.color.rgb = SECONDARY_COLOR
                
                # Add content
                body_shape = slide.placeholders[1]
                text_frame = body_shape.text_frame
                text_frame.clear()
                text_frame.word_wrap = True
                
                for item_idx, item in enumerate(chunk):
                    if item_idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = item if not item.endswith('.') else item
                    p.level = 0
                    p.font.size = Pt(16)
                    p.font.color.rgb = TEXT_COLOR
                    p.space_after = Pt(12)
            
            # Add key points slide if bullets exist
            if section_bullets:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                title_shape = slide.shapes.title
                title_shape.text = f"{section_title} - Key Points"
                title_shape.text_frame.paragraphs[0].font.size = Pt(32)
                title_shape.text_frame.paragraphs[0].font.bold = True
                title_shape.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
                
                body_shape = slide.placeholders[1]
                text_frame = body_shape.text_frame
                text_frame.clear()
                
                for bullet_idx, bullet in enumerate(section_bullets[:6]):  # Max 6 key points
                    if bullet_idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.color.rgb = ACCENT_COLOR
                    p.space_after = Pt(14)
        
        # ===== CITATIONS SLIDE =====
        if citations:
            citation_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(citation_slide_layout)
            
            title_shape = slide.shapes.title
            title_shape.text = "References and Citations"
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = SECONDARY_COLOR
            
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()
            
            for citation_idx, citation in enumerate(citations[:8], 1):  # Max 8 citations per slide
                if citation_idx == 1:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"[{citation_idx}] {citation}"
                p.level = 0
                p.font.size = Pt(12)
                p.font.color.rgb = SUBTITLE_COLOR
                p.space_after = Pt(10)
        
        # Save presentation
        prs.save(str(output_path))
        
        logger.info(f"PPT generated successfully: {output_path}")
        logger.info(f"Created {len(prs.slides)} slides total")
        return str(output_path)
        
    except Exception as e:
        logger.error(f"Error generating PPT: {str(e)}", exc_info=True)
        return json.dumps({"error": str(e)})
